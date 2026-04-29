"""Generate PowerPoint decks for P4 paper portfolio.

Outputs:
  papers/P4_llm_as_judge/slides/P4-SHORT.pptx (~14 slides, LLM4Eval workshop, 15-20 min)
  papers/P4_llm_as_judge/slides/P4-LONG.pptx  (~32 slides, ECIR full conference, 25-30 min)

Re-run anytime to regenerate from the latest figures + numbers in this script.

Usage:
  py -3 -X utf8 papers/P4_llm_as_judge/slides/make_p4_slides.py
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent
FIGS = ROOT.parent / "figures"
OUT_SHORT = ROOT / "P4-SHORT.pptx"
OUT_LONG = ROOT / "P4-LONG.pptx"

# ----------------------- Visual theme -----------------------
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_TITLE = RGBColor(0x0E, 0x4A, 0x86)        # ISU cardinal-ish dark blue
COLOR_ACCENT = RGBColor(0xC8, 0x10, 0x2E)        # ISU cardinal red
COLOR_GREEN = RGBColor(0x2E, 0x8B, 0x57)         # reasoning-generous cluster
COLOR_BLUE = RGBColor(0x4A, 0x6F, 0xA5)          # strict-mid cluster
COLOR_GREY = RGBColor(0x70, 0x70, 0x70)
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"


# ----------------------- Helpers -----------------------
def new_pres(title_bg_color=None) -> Presentation:
    """Return a fresh 16:9 Presentation."""
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


def add_blank(p: Presentation):
    """Add a blank slide and return its slide object."""
    blank_layout = p.slide_layouts[6]
    return p.slides.add_slide(blank_layout)


def add_textbox(slide, left, top, width, height, text, *,
                size=18, bold=False, color=COLOR_TEXT,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY):
    """Add a text box with single-paragraph text."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=18, color=COLOR_TEXT, font=FONT_BODY, line_spacing=1.15):
    """Add a bulleted list. items is list of strings (or list of (text, bold) tuples)."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, bold = item
        else:
            text, bold = item, False
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.line_spacing = line_spacing
        run = para.add_run()
        run.text = "• " + text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_title(slide, text, *, size=32, color=COLOR_TITLE, top=0.35):
    """Standard slide title at the top, single line."""
    return add_textbox(slide, 0.5, top, 12.3, 0.7, text,
                       size=size, bold=True, color=color, align=PP_ALIGN.LEFT, font=FONT_TITLE)


def add_subtitle(slide, text, *, size=16, color=COLOR_GREY, top=1.05):
    return add_textbox(slide, 0.5, top, 12.3, 0.5, text,
                       size=size, bold=False, color=color, align=PP_ALIGN.LEFT, font=FONT_BODY)


def add_image(slide, path: Path, left, top, width=None, height=None):
    """Insert an image. If only width or height given, the other is auto-scaled."""
    if not path.exists():
        # Placeholder if image missing
        add_textbox(slide, left, top, width or 6, height or 4,
                    f"[missing: {path.name}]", size=14, color=COLOR_ACCENT)
        return
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


def add_table(slide, left, top, rows_data, *,
              col_widths=None, header_bold=True, font_size=14,
              header_bg=COLOR_TITLE, header_fg=RGBColor(0xFF, 0xFF, 0xFF),
              row_bg=None, row_alt=COLOR_LIGHT_BG, total_width=None):
    """Add a table from a 2D list. First row = header.
    rows_data: list of lists of strings.
    col_widths: list of float (inches), summing to total_width if given.
    """
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    if total_width is None:
        total_width = 12.0
    if col_widths is None:
        col_widths = [total_width / n_cols] * n_cols

    height_per_row = 0.34
    total_height = height_per_row * n_rows

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(total_width), Inches(total_height),
    )
    tbl = tbl_shape.table

    # Set column widths
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            run = para.add_run()
            run.text = str(cell_text)
            run.font.size = Pt(font_size)
            run.font.name = FONT_BODY
            if ri == 0:
                run.font.bold = header_bold
                run.font.color.rgb = header_fg
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                run.font.color.rgb = COLOR_TEXT
                cell.fill.solid()
                if row_alt is not None and ri % 2 == 0:
                    cell.fill.fore_color.rgb = row_alt
                else:
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
    return tbl_shape


def add_footer(slide, text, *, size=10, color=COLOR_GREY):
    return add_textbox(slide, 0.5, 7.05, 12.3, 0.3, text,
                       size=size, color=color, align=PP_ALIGN.LEFT)


def add_slide_number(slide, n, total):
    return add_textbox(slide, 12.5, 7.05, 0.7, 0.3, f"{n}/{total}",
                       size=10, color=COLOR_GREY, align=PP_ALIGN.RIGHT)


# ----------------------- Common slide patterns -----------------------
def title_slide(p, title, subtitle, authors_line):
    s = add_blank(p)
    # Centered title
    add_textbox(s, 1, 2.4, 11.3, 1.2, title,
                size=40, bold=True, color=COLOR_TITLE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
    add_textbox(s, 1, 3.7, 11.3, 0.6, subtitle,
                size=22, color=COLOR_GREY, align=PP_ALIGN.CENTER)
    add_textbox(s, 1, 4.5, 11.3, 0.5, authors_line,
                size=18, color=COLOR_TEXT, align=PP_ALIGN.CENTER, bold=True)
    add_textbox(s, 1, 5.05, 11.3, 0.4, "Iowa State University",
                size=16, color=COLOR_GREY, align=PP_ALIGN.CENTER)
    # Bottom strip with venue
    add_textbox(s, 1, 6.7, 11.3, 0.4, "asukul@iastate.edu  ·  2026-04-28",
                size=12, color=COLOR_GREY, align=PP_ALIGN.CENTER)
    return s


def section_divider(p, label, big_text):
    s = add_blank(p)
    # Color band on left
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_ACCENT
    band.line.fill.background()
    add_textbox(s, 1.2, 2.7, 11, 0.5, label.upper(),
                size=22, color=COLOR_GREY, bold=True, align=PP_ALIGN.LEFT)
    add_textbox(s, 1.2, 3.3, 11, 1.5, big_text,
                size=44, color=COLOR_TITLE, bold=True, align=PP_ALIGN.LEFT, font=FONT_TITLE)
    return s


def standard_content_slide(p, title, subtitle=None):
    s = add_blank(p)
    add_title(s, title)
    if subtitle:
        add_subtitle(s, subtitle)
    return s


# ----------------------- Slide generators -----------------------
# Each function takes (presentation, slide_num, total) and returns nothing.
# Slide number tracker is updated externally.

def slide_title(p, n, total, *, version_label):
    s = title_slide(
        p,
        title="Cross-Family LLM-Judge Agreement\nfor Institutional RAG",
        subtitle="A 5-Family, 9-Judge Ablation with Mechanism + External Validation",
        authors_line="Adisak Sukul",
    )
    add_slide_number(s, n, total)
    add_footer(s, f"P4 — {version_label}")


def slide_motivation(p, n, total, *, version_label):
    s = standard_content_slide(p, "The which-judge problem")
    add_subtitle(s, "RAG over institutional corpora has no gold relevance set; LLM-as-Judge fills the gap — but which judge?")
    add_bullets(s, 0.7, 1.7, 12, 4.8, [
        ("No gold qrels for institutional corpora — TREC-style ground truth needs hundreds of expert hours per topic", True),
        ("Generic IR benchmarks (BEIR, MS-MARCO) don't transfer — web-style queries on web-style corpora", False),
        ("LLM-as-Judge surfaces a new question: WHICH judge? Sonnet vs GPT-5 vs Gemini vs Qwen vs Gemma?", True),
        ("Prior work compares within ONE provider family, or two-family pairs without ordinal weighting, or excludes open-weight peers entirely", False),
        ("Open-weight judges are the default for sovereign-cloud and privacy-conservative institutions — excluding them silently endorses commercial-only practice", True),
    ], size=18, line_spacing=1.3)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_contribution_overview(p, n, total, *, version_label):
    s = standard_content_slide(p, "Four contribution claims")
    add_subtitle(s, "First multi-family pairwise κ matrix with within-family controls AND open-weight as first-class judges")
    add_bullets(s, 0.7, 1.7, 12, 4.8, [
        ("C1: Cross-family reasoning judges converge at κ ≥ 0.75 — well above the 0.4-0.6 typical baselines", True),
        ("C2: Within-family agreement is BOUNDED by the cross-family ceiling (no within-family pair dominates)", True),
        ("C3: Open-weight judges (Qwen × Gemma 4) produce the matrix-highest within-pair κ = 0.80 at ~1% of commercial cost", True),
        ("C4: Open-source toolkit (eval_llm_judge.py, validate_against_trec.py) + community disclosure template", True),
        ("Plus: §6 mechanistic decomposition (R²=93% from joint distribution; tokenizer hypothesis refuted)", False),
        ("Plus: §7 external validation against NIST TREC RAG 2024 qrels (ensemble κ = 0.4941 on 537 stratified-balanced pairs)", False),
    ], size=17, line_spacing=1.25)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_judge_slate(p, n, total, *, version_label):
    s = standard_content_slide(p, "Judge slate: 9 judges across 5 families, 4 within-family pairs")
    rows = [
        ["#", "Judge", "Family", "Reasoning", "Route"],
        ["1", "Claude Opus 4.7", "Anthropic", "yes", "OpenRouter"],
        ["2", "Claude Sonnet 4.6", "Anthropic", "yes", "OpenRouter"],
        ["3", "GPT-5.5 (reasoning=low)", "OpenAI", "yes", "OpenAI direct"],
        ["4", "GPT-4o (chat)", "OpenAI", "no", "OpenAI direct"],
        ["5", "Gemini 3.1 Pro Preview", "Google-commercial", "yes (thinking)", "OpenRouter"],
        ["6", "Gemini 2.5 Pro", "Google-commercial", "yes (thinking)", "OpenRouter"],
        ["7", "DeepSeek V4 Pro", "DeepSeek", "yes", "OpenRouter"],
        ["8", "Qwen 3.6 Plus", "Open-weight", "no", "OpenRouter"],
        ["9", "Gemma 4 26B", "Open-weight", "no", "OpenRouter"],
    ]
    add_table(s, 0.7, 1.7, rows,
              col_widths=[0.4, 3.0, 2.2, 1.7, 2.0],
              total_width=9.3, font_size=14)
    add_textbox(s, 0.7, 5.7, 12, 1.3,
                "Within-family pairs: Anthropic (1↔2), OpenAI (3↔4), Google-commercial (5↔6), Open-weight cross-organization (8↔9). "
                "All 9 judges score the same 570 (query, document) pairs from ISU DSpace (97k full-text PDFs) — "
                "score variance is purely judge-level. Versions pinned 2026-04-25.",
                size=13, color=COLOR_GREY)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_pipeline_diagram(p, n, total, *, version_label):
    s = standard_content_slide(p, "Pipeline: 9-judge ThreadPool fan-out")
    add_subtitle(s, "eval_llm_judge.py — multi-judge mode, retry on transient errors, JSON output per judge")
    add_image(s, FIGS / "diagram_3_our_multijudge_fanout.png", 0.7, 1.7, width=12)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_per_judge_metrics(p, n, total, *, version_label):
    s = standard_content_slide(p, "Per-judge retrieval metrics — same documents, different verdicts")
    rows = [
        ["Judge", "Family", "nDCG@10", "P@5", "MRR", "Mean"],
        ["Sonnet 4.6", "Anthropic", "0.862", "0.649", "0.826", "1.68"],
        ["GPT-5.5 (low)", "OpenAI", "0.846", "0.600", "0.792", "1.63"],
        ["DSV4 Pro", "DeepSeek", "0.835", "0.702", "0.883", "1.63"],
        ["GPT-4o", "OpenAI", "0.803", "0.361", "0.575", "1.15"],
        ["Qwen 3.6", "Open-weight", "0.758", "0.379", "0.599", "1.18"],
        ["Gemma 4 26B", "Open-weight", "0.753", "0.375", "0.573", "1.11"],
        ["Opus 4.7", "Anthropic", "0.749", "0.372", "0.536", "1.09"],
        ["Gemini 2.5 Pro", "Google", "0.624", "0.281", "0.486", "0.76"],
        ["Gemini 3.1 Prev", "Google", "0.454", "0.158", "0.315", "0.37"],
    ]
    add_table(s, 0.7, 1.5, rows,
              col_widths=[2.2, 2.0, 1.4, 1.2, 1.2, 1.0],
              total_width=9.0, font_size=13)
    add_textbox(s, 0.7, 5.6, 12, 1.4,
                "Same retrieved documents → nDCG@10 between 0.45 and 0.86 — a 1.9× spread driven entirely by judge selection. "
                "Two practitioners running identical retrieval code, scoring with different flagship judges, can arrive at "
                "OPPOSITE go/no-go decisions. Disclosure template (C4): \"nDCG@10 = X.XX via [family]/[model]/[reasoning], N pairs, DATE\".",
                size=13, color=COLOR_GREY)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_kappa_matrix_9j(p, n, total, *, version_label):
    s = standard_content_slide(p, "9×9 pairwise κ matrix — every off-diagonal ≥ 0.56")
    add_subtitle(s, "Quadratic-weighted Cohen's κ on 570 pairs · two clusters · cross-org open-weight = 0.80 (matrix-highest)")
    add_image(s, FIGS / "kappa_matrix_9judge.png", 2.5, 1.6, height=5.4)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_calibration_topology(p, n, total, *, version_label):
    s = standard_content_slide(p, "Calibration clusters — calibration philosophy > provider lineage")
    add_image(s, FIGS / "diagram_5_calibration_topology.png", 0.7, 1.5, width=12)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_C1(p, n, total, *, version_label):
    s = standard_content_slide(p, "C1: Cross-family reasoning judges converge at κ ≥ 0.75")
    add_bullets(s, 0.7, 1.4, 12.5, 5.0, [
        ("Five cross-family pairs reach substantial-or-better κ (Sonnet ↔ GPT-5.5 = 0.79, Qwen ↔ Gemma 4 = 0.80, etc.)", False),
        ("Well above the 0.4-0.6 typical of prior art [Rahmani 2024] and the 0.60 unweighted baseline of [Thakur 2025]", False),
        ("DeepSeek V4 Pro joining the cluster RULES OUT a Western-training-data common-cause explanation", True),
        ("Pattern replicates across three independent ablation scales (5-, 7-, 9-judge) within run-to-run noise", False),
        ("EXTERNAL CORROBORATION: 9-judge ensemble achieves κ = 0.49 against published NIST qrels on TREC RAG 2024 — "
         "the convergence is not specific to our institutional corpus", True),
    ], size=18, line_spacing=1.35)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_C2(p, n, total, *, version_label):
    s = standard_content_slide(p, "C2: Within-family agreement is BOUNDED by the cross-family ceiling")
    rows = [
        ["Within-family pair", "κ", "vs cross-family ceiling (0.79)"],
        ["Anthropic (Opus ↔ Sonnet)", "0.71", "below"],
        ["OpenAI (GPT-5.5 ↔ GPT-4o)", "0.63", "below"],
        ["Google-commercial (Gemini 3.1 ↔ 2.5)", "0.67", "below"],
        ["Cross-organization Open-weight (Qwen ↔ Gemma 4)", "0.80", "MATRIX MAX"],
    ]
    add_table(s, 0.7, 1.5, rows, col_widths=[6.0, 1.0, 4.0],
              total_width=11.0, font_size=15)
    add_textbox(s, 0.7, 4.5, 12, 2.5,
                "No within-family pair DOMINATES — at odds with self-preference findings on open-ended generation [Panickssery 2024]. "
                "Our companion 4-way extraction study reports within-family agreement ~2× cross-family on open-vocabulary tasks; "
                "the inversion under bounded ordinal judging suggests SELF-PREFERENCE IS MEDIATED BY OUTPUT-SPACE BOUNDEDNESS, "
                "not provider lineage. In a 4-class output space, the dominant signal is calibration philosophy, which crosses provider lines.",
                size=15, color=COLOR_TEXT)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_C3(p, n, total, *, version_label):
    s = standard_content_slide(p, "C3: Cross-org open-weight = matrix-highest κ at ~1% of commercial cost")
    add_bullets(s, 0.7, 1.4, 12.5, 5.0, [
        ("Qwen 3.6 Plus ↔ Gemma 4 26B = κ 0.80 — exceeds every commercial within-family pair", True),
        ("Both calibrate with strict-mid commercial models (mean 1.11-1.18), NOT with reasoning-generous (1.63-1.68)", False),
        ("Marginal cost ~USD 0.30 (vs ~USD 18 for the commercial 7-judge frontier)", True),
        ("100% coverage on TREC RAG 2024 — better than 4 of 7 frontier judges (Gemini 17-24%, DSV4 40%)", True),
        ("Direct policy implication: sovereign-cloud and privacy-conservative deployments are NOT condemned to lower-quality judging — "
         "the cross-org open-weight pair is the recommended low-cost configuration", True),
    ], size=18, line_spacing=1.35)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_C4(p, n, total, *, version_label):
    s = standard_content_slide(p, "C4: Open-source toolkit + community disclosure template")
    add_bullets(s, 0.7, 1.4, 12.5, 5.0, [
        ("eval_llm_judge.py — multi-judge mode, ThreadPool fan-out, p4-frontier + p4-supplement-openweight presets", False),
        ("validate_against_trec.py — multi-corpus external-validation harness (TREC RAG 2024, BEIR scifact, TREC-COVID)", False),
        ("fetch_msmarco_passages.py — HF-streaming MS MARCO v2.1 passage extractor (peak disk ~400 MB, 38 min wall)", False),
        ("All 9 within-corpus per-judge JSONs (5,130 records) + 9 external-validation JSONs (4,833 records)", False),
        ("Cross-run merge harness with retrieval-determinism verification + κ matrix regen scripts", False),
        ("Disclosure norm: nDCG@10 = X.XX via [family]/[model]/[reasoning], N pairs, DATE — κ heatmap with every nDCG plot", True),
    ], size=18, line_spacing=1.35)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_mechanism_summary(p, n, total, *, version_label):
    s = standard_content_slide(p, "Mechanism: joint-distribution structure explains R² = 93% of κ variance")
    add_image(s, FIGS / "judge_calibration_mechanism.png", 0.5, 1.3, width=12.3)
    add_textbox(s, 0.5, 6.4, 12.3, 0.6,
                "Marginal KL = coarse 33% proxy · joint-distribution dispersion + effective rank = R²=93% · "
                "structural factors fully mediated · shared-tokenizer hypothesis REFUTED (Qwen↔Gemma Jaccard=0.066, lowest, yet κ=0.80 highest).",
                size=12, color=COLOR_GREY)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_external_validation_setup(p, n, total, *, version_label):
    s = standard_content_slide(p, "External validation: NIST TREC RAG 2024 — 537 stratified-balanced pairs")
    add_bullets(s, 0.7, 1.4, 12.5, 5.0, [
        ("Sample: 537 pairs from official 20,283-pair NIST qrels (random.seed(42))", False),
        ("Stratified-balanced: 135 / 134 / 134 / 134 across labels 0/1/2/3 — uniform over the 4 ordinal classes", True),
        ("Coverage: all 86 unique queries; 537 unique passages from MS MARCO v2.1 segmented corpus", False),
        ("Passage extraction: HF-stream from drexalt/msmarco-2.1-segmented (60 shards, ~400 MB peak disk, 38 min wall)", False),
        ("Cost: ~USD 35, ~7.5 h wall (vs IRB human study estimate $1,500-3,000, 6-12 weeks)", True),
        ("BEIR scifact also run (300 pairs, all-positive qrels — κ undefined; reports precision instead = 63.7% ensemble)", False),
    ], size=18, line_spacing=1.35)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_external_validation_results(p, n, total, *, version_label):
    s = standard_content_slide(p, "External validation results — TREC RAG 2024 9-judge κ vs NIST qrels")
    rows = [
        ["Rank", "Judge", "κ", "valid/537"],
        ["1", "Gemini 2.5 Pro", "0.5513", "92 (small-n)"],
        ["2", "Claude Sonnet 4.6", "0.5123", "537"],
        ["3", "Claude Opus 4.7", "0.4792", "537"],
        ["4", "GPT-5.5 (reasoning=low)", "0.4789", "537"],
        ["5", "DeepSeek V4 Pro", "0.4705", "212"],
        ["6", "Qwen 3.6 Plus", "0.4141", "537"],
        ["7", "Gemini 3.1 Pro Preview", "0.4092", "127"],
        ["8", "GPT-4o (chat)", "0.4065", "537"],
        ["9", "Gemma 4 26B", "0.3958", "537"],
        ["", "9-judge ensemble median", "0.4941", "537"],
        ["", "7-judge frontier-only median", "0.5187", "537"],
    ]
    add_table(s, 0.7, 1.4, rows,
              col_widths=[0.6, 4.0, 1.2, 2.0],
              total_width=7.8, font_size=14)
    add_textbox(s, 9.0, 1.4, 4.0, 5.0,
                "Headlines\n\n"
                "9-judge ensemble κ = 0.4941\n→ moderate, near-substantial\n\n"
                "5 of 9 judges hit κ ≥ 0.47\n\n"
                "All 9 judges κ > 0.39 — none at chance\n\n"
                "Ensemble outperforms 7 of 9 single judges\n\n"
                "Adding open-weight broadens coverage,\nslightly lowers κ (robustness/headline tradeoff)",
                size=13, color=COLOR_TEXT)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_coverage_finding(p, n, total, *, version_label):
    s = standard_content_slide(p, "Coverage divergence — content-domain reliability axis")
    rows = [
        ["Judge", "TREC RAG 2024 (n=537)", "BEIR scifact (n=300)"],
        ["GPT-5.5 (reasoning=low)", "100%", "100%"],
        ["GPT-4o", "100%", "100%"],
        ["Claude Opus 4.7", "100%", "95%"],
        ["Claude Sonnet 4.6", "100%", "99.7%"],
        ["Qwen 3.6 Plus", "100%", "100%"],
        ["Gemma 4 26B", "100%", "100%"],
        ["DeepSeek V4 Pro", "39%", "84%"],
        ["Gemini 3.1 Pro Preview", "24%", "60%"],
        ["Gemini 2.5 Pro", "17%", "86%"],
    ]
    add_table(s, 0.7, 1.5, rows,
              col_widths=[3.5, 3.0, 3.0],
              total_width=9.5, font_size=14)
    add_textbox(s, 0.7, 6.0, 12, 1.0,
                "ALWAYS-WORKS 6-judge subset (≥95% on both): Anthropic + OpenAI + Qwen + Gemma. "
                "The 2 cheapest open-weight judges make this subset. "
                "Gemini fails on TREC RAG (web content) more than on BEIR (scientific abstracts) — reliability is content-domain dependent.",
                size=14, color=COLOR_TEXT, bold=False)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_takeaway(p, n, total, *, version_label):
    s = standard_content_slide(p, "Takeaway")
    add_bullets(s, 0.7, 1.4, 12.5, 5.0, [
        ("CALIBRATION PHILOSOPHY (lenient vs strict score-allocation pattern), NOT provider lineage, "
         "is the axis on which retrieval verdicts shift", True),
        ("Cross-family agreement (κ ≥ 0.75) is high enough that judge interchangeability is a defensible default "
         "for reasoning-capable commercial judges", False),
        ("Cross-organization open-weight ensemble (Qwen + Gemma 4) achieves matrix-highest κ at ~1% commercial cost — "
         "sovereign-cloud / on-prem deployments are NOT condemned to lower-quality judging", True),
        ("External NIST validation (κ = 0.49 ensemble) confirms moderate-substantial agreement on a public IR benchmark", False),
        ("Mechanism: 93% of pairwise κ variance is joint-distribution structure of paired scores — a NEW insight for ensemble design", True),
    ], size=18, line_spacing=1.35)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_qa(p, n, total, *, version_label):
    s = add_blank(p)
    add_textbox(s, 1, 2.5, 11.3, 1.5, "Q&A",
                size=80, bold=True, color=COLOR_TITLE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
    add_textbox(s, 1, 4.2, 11.3, 0.6, "asukul@iastate.edu",
                size=22, color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    add_textbox(s, 1, 4.8, 11.3, 0.5,
                "Code + data + paper draft: github.com/asukul/RAG-Eval-LLM-Judge",
                size=18, color=COLOR_GREY, align=PP_ALIGN.CENTER)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


# ----------------------- Long-version-specific slides -----------------------

def slide_outline_long(p, n, total):
    s = standard_content_slide(p, "Outline")
    add_bullets(s, 1.0, 1.4, 12, 5.5, [
        ("§1 Introduction — the which-judge problem", True),
        ("§2 Related work — 2025 inter-LLM-judge cluster", False),
        ("§3 Methodology — corpus, 9 judges, rubric, pipeline", False),
        ("§4 Within-corpus results — metrics, calibration tiers, κ matrix", True),
        ("§5 Findings — C1, C2, C3, C4", True),
        ("§6 Mechanism — joint distribution R²=93%, mediation, tokenizer refutation", True),
        ("§7 External validation — NIST TREC RAG 2024 + BEIR scifact + coverage", True),
        ("§8 Limitations · §9 Reproducibility · §10 Conclusion", False),
    ], size=20, line_spacing=1.4)
    add_footer(s, "P4 — LONG")
    add_slide_number(s, n, total)


def slide_related_work(p, n, total):
    s = standard_content_slide(p, "Related work — 2025 inter-LLM-judge agreement cluster")
    rows = [
        ["Work", "Families covered", "Within-family pairs", "Open-weight as 1st-class", "κ metric"],
        ["UMBRELA [2024]", "1 (deployed)", "no", "no", "%agreement"],
        ["Balog 2025", "1", "yes", "no", "Cohen κ"],
        ["Farzi 2025", "1", "no", "no", "Cohen κ"],
        ["Thakur 2025", "2 (GPT-4o, Llama)", "no", "partial", "unweighted κ"],
        ["Han 2025", "judge↔human only", "—", "—", "Cohen κ"],
        ["Ours", "5", "4 controls", "yes (Qwen × Gemma)", "qw-Cohen κ"],
    ]
    add_table(s, 0.5, 1.5, rows,
              col_widths=[2.5, 2.2, 2.0, 3.0, 2.5],
              total_width=12.2, font_size=12)
    add_textbox(s, 0.5, 5.5, 12, 1.4,
                "GAP: no published paper, to our knowledge, reports a multi-family pairwise κ matrix that includes "
                "BOTH within-family controls AND open-weight peers as first-class judges at ≥ 500 quadratic-weighted ordinal observations. "
                "P4 closes that gap with 5 families × 4 controls × 9 judges × 570 pairs (within) + 537 pairs (NIST external).",
                size=14, color=COLOR_TEXT)
    add_footer(s, "P4 — LONG")
    add_slide_number(s, n, total)


def slide_corpus_setup(p, n, total):
    s = standard_content_slide(p, "Corpus and query set — Iowa State University DSpace")
    add_bullets(s, 0.7, 1.4, 12.5, 5.0, [
        ("ISU DSpace: 97,441 full-text PDFs → 1.03M chunks (≤500 words, 100-word overlap)", False),
        ("Embedding: Google Vertex text-embedding-005 in Qdrant 1.13 HNSW (M=16, ef_construct=100, cosine)", False),
        ("Query set: REFINE-synthesized 57 queries [Sukul, forthcoming]", True),
        ("    Intents: factoid / methodological / comparative / author-style — balanced across 5 intents", False),
        ("    Coverage: 12:11:12:11:11 across 20 topic clusters", False),
        ("Top-10 retrieval per query → 570 (query, document) pairs scored by all 9 judges", True),
        ("Score variance is purely judge-level — same documents, same order, no retrieval randomness", False),
    ], size=18, line_spacing=1.3)
    add_footer(s, "P4 — LONG")
    add_slide_number(s, n, total)


def slide_rubric(p, n, total):
    s = standard_content_slide(p, "Rubric — bounded 0-3 ordinal scale")
    rows = [
        ["Score", "Definition", "Example trigger"],
        ["0", "Irrelevant — document does not address the query at all", "Topic mismatch"],
        ["1", "Topically related — mentions the topic but does not answer", "Background context only"],
        ["2", "Partially answers — provides some of what the query asks for", "Direct mention but incomplete"],
        ["3", "Fully answers — directly answers the query", "Specific, complete information"],
    ]
    add_table(s, 0.7, 1.4, rows, col_widths=[0.8, 6.5, 4.7],
              total_width=12.0, font_size=15)
    add_textbox(s, 0.7, 4.6, 12, 2.0,
                "Calibrated at prompt time with FIVE worked examples drawn from a held-out seed set. "
                "Documents truncated to 1,500 chars to control verbosity bias [Saito 2023]. "
                "Quadratic weighting on disagreements: a 0↔3 disagreement counts 9× a 0↔1 disagreement — "
                "captures the ordinal structure that simple percent-agreement loses.",
                size=14, color=COLOR_TEXT)
    add_footer(s, "P4 — LONG")
    add_slide_number(s, n, total)


def slide_pipeline_diagram_long(p, n, total):
    s = standard_content_slide(p, "Multi-judge pipeline — ThreadPool fan-out vs standard single-judge")
    add_image(s, FIGS / "diagram_4_norm_vs_ours.png", 0.5, 1.4, width=12.3)
    add_footer(s, "P4 — LONG")
    add_slide_number(s, n, total)


def slide_calibration_tiers(p, n, total):
    s = standard_content_slide(p, "Three calibration tiers — under None=0 vs valid-only accounting")
    rows = [
        ["Judge", "Mean (None=0)", "Mean (valid)", "Tier (valid-only)"],
        ["Sonnet 4.6", "1.68", "1.72", "Reasoning-generous"],
        ["GPT-5.5 (low)", "1.63", "1.63", "Reasoning-generous"],
        ["DSV4 Pro", "1.63", "1.67", "Reasoning-generous"],
        ["Opus 4.7", "1.09", "1.26", "Strict-mid"],
        ["GPT-4o", "1.15", "1.15", "Strict-mid"],
        ["Qwen 3.6", "1.18", "1.18", "Strict-mid"],
        ["Gemma 4 26B", "1.11", "1.11", "Strict-mid"],
        ["Gemini 2.5 Pro", "0.76", "1.45", "Transitional → strict-mid (revised)"],
        ["Gemini 3.1 Prev", "0.37", "1.17", "Strict-mid (revised, was outlier)"],
    ]
    add_table(s, 0.4, 1.4, rows,
              col_widths=[2.4, 2.0, 2.0, 5.4],
              total_width=11.8, font_size=14)
    add_textbox(s, 0.4, 5.5, 12, 1.4,
                "STRICT-OUTLIER tier dissolves under valid-only accounting — Gemini judges' apparent strictness was a missing-data artifact "
                "(Gemini 3.1 Prev: 68.4% None rate; 2.5 Pro: 47.9%). Both open-weight judges (mean 1.11-1.18) calibrate with strict-mid "
                "commercial models (GPT-4o, Opus 4.7) — non-obvious given the training-compute gap.",
                size=12, color=COLOR_GREY)
    add_footer(s, "P4 — LONG")
    add_slide_number(s, n, total)


def slide_pair_confusion_top(p, n, total):
    s = standard_content_slide(p, "§6 — 4 highest-κ pairs (joint distribution: tight diagonal)")
    add_image(s, FIGS / "judge_pair_confusion_top4_highest.png", 1.2, 1.6, width=11.0)
    add_textbox(s, 1.2, 4.7, 11, 1.8,
                "Top-4 κ pairs (Qwen↔Gemma 4 = 0.797, Sonnet↔GPT-5.5 = 0.789, GPT-5.5↔Gem 2.5 = 0.776, Gemma 4↔GPT-4o = 0.771). "
                "Each 4×4 confusion matrix shows P(judge A score = i, judge B score = j); diagonal concentration means "
                "both judges agree on which score each document deserves. The dispersion (quadratic-weighted disagreement distance) "
                "of these matrices is the dominant predictor of pairwise κ.",
                size=14, color=COLOR_TEXT)
    add_footer(s, "P4 — LONG")
    add_slide_number(s, n, total)


def slide_pair_confusion_bottom(p, n, total):
    s = standard_content_slide(p, "§6 — 4 lowest-κ pairs (joint distribution: off-diagonal dispersion)")
    add_image(s, FIGS / "judge_pair_confusion_bottom4_lowest.png", 1.2, 1.6, width=11.0)
    add_textbox(s, 1.2, 4.7, 11, 1.8,
                "Bottom-4 κ pairs (Sonnet↔Gem 3.1 Prev = 0.560, Gem 3.1 Prev↔DSV4 Pro = 0.573, Sonnet↔Gemma 4 = 0.620, GPT-5.5↔GPT-4o = 0.632). "
                "Off-diagonal mass means the judges disagree by 1-2 score levels on the same documents — "
                "tier-crossing pairs (reasoning-generous vs strict-mid) systematically lower in κ. "
                "Together TOP and BOTTOM rows: dispersion + effective rank → R² = 93% of pairwise κ variance explained.",
                size=14, color=COLOR_TEXT)
    add_footer(s, "P4 — LONG")
    add_slide_number(s, n, total)


def slide_regression_AB(p, n, total):
    s = standard_content_slide(p, "§6 — Regression decomposition (A: cumulative R²; B: incremental ΔR²)")
    add_image(s, FIGS / "judge_kappa_regression_panels_AB.png", 1.5, 1.5, width=10.5)
    add_textbox(s, 1.5, 5.7, 10.5, 1.4,
                "Panel A: nested OLS models — M1 (κ ~ KL) reaches R²=0.326; adding reasoning-mode (M2: 0.342), "
                "provider (M3: 0.354), model class (M4: 0.384) yields only marginal improvements. "
                "Panel B: KL alone contributes ΔR² = 0.326; the three structural predictors add only 0.015, 0.013, 0.030 — "
                "marginal-distribution similarity is the dominant explainable factor among single-feature models.",
                size=14, color=COLOR_TEXT)
    add_footer(s, "P4 — LONG")
    add_slide_number(s, n, total)


def slide_regression_CD(p, n, total):
    s = standard_content_slide(p, "§6 — Regression decomposition (C: M4 coefficients; D: predicted vs actual)")
    add_image(s, FIGS / "judge_kappa_regression_panels_CD.png", 1.5, 1.5, width=10.5)
    add_textbox(s, 1.5, 5.7, 10.5, 1.4,
                "Panel C: full-model coefficients with 95% CI — KL is the only significant predictor (β=-0.130, p<0.001); "
                "same_reasoning, same_provider, same_class all p > 0.18 (gray = non-significant). "
                "Panel D: M4-predicted κ vs actual κ — points cluster near the y=x line; full R² jumps to 0.93 "
                "when joint-distribution features (dispersion + effective rank) replace KL as input — see companion long version §6.2.",
                size=14, color=COLOR_TEXT)
    add_footer(s, "P4 — LONG")
    add_slide_number(s, n, total)


def slide_tokenizer_refutation(p, n, total):
    s = standard_content_slide(p, "§6 mechanism — shared-tokenizer hypothesis REFUTED")
    rows = [
        ["Pair", "Tokenizer Jaccard", "κ", "Verdict"],
        ["Anthropic Opus ↔ Sonnet", "1.000", "0.71", "Same family, mid κ"],
        ["Gemini 3.1 ↔ Gemini 2.5", "0.993", "0.67", "Same family, mid κ"],
        ["GPT-5.5 ↔ GPT-4o", "≈ 0.99", "0.63", "Same family, lowest in-family κ"],
        ["Qwen 3.6 ↔ Gemma 4", "0.066 (lowest in slate!)", "0.80 (matrix max!)", "REFUTES the hypothesis"],
    ]
    add_table(s, 0.4, 1.4, rows,
              col_widths=[3.5, 2.8, 2.5, 3.0],
              total_width=11.8, font_size=14)
    add_textbox(s, 0.5, 4.4, 12, 2.0,
                "If shared tokenizer drove κ, Qwen↔Gemma 4 should be at the BOTTOM of the κ ranking — instead they are at the TOP. "
                "The convergence of Qwen and Gemma 4 happens at the DECISION-MAKING LAYER, not the input-encoding layer. "
                "Pairwise κ is not driven by shared input vocabulary; the calibration cluster effect comes from how each judge "
                "maps inputs to scores after encoding.",
                size=15, color=COLOR_TEXT)
    add_footer(s, "P4 — LONG")
    add_slide_number(s, n, total)


def slide_thakur_comparison(p, n, total):
    s = standard_content_slide(p, "Comparison to Thakur 2025 baseline")
    add_bullets(s, 0.7, 1.4, 12.5, 5.0, [
        ("Thakur et al. 2025: GPT-4o single judge, κ = 0.60 on 537-pair \"support evaluation\" subset", True),
        ("Our GPT-4o on our 537-pair stratified-balanced sample: κ = 0.41 (gap of 0.19)", True),
        ("Our 9-judge ensemble: κ = 0.49 — partially closes the gap (+0.08 over our GPT-4o)", False),
        ("Likely explanations for the gap (in decreasing impact order):", True),
        ("    Sample selection — ours stratified-balanced (uniform), theirs natural-distribution-proportional", False),
        ("    Prompt template — exact Thakur prompt not specified in their paper", False),
        ("    Qrels version — possible UMBRELA-augmented vs official NIST", False),
        ("    Judge model version — we pin gpt-4o-2024-08-06; theirs unspecified", False),
        ("Direct apples-to-apples replication awaits receipt of their curated 537-pair subset", False),
    ], size=16, line_spacing=1.3)
    add_footer(s, "P4 — LONG")
    add_slide_number(s, n, total)


def slide_limitations(p, n, total, *, version_label):
    s = standard_content_slide(p, "Limitations")
    add_bullets(s, 0.7, 1.4, 12.5, 5.4, [
        ("Single internal corpus (ISU DSpace) — open-weight κ ceiling most likely to drift across domains", False),
        ("Single internal query set (REFINE 57 queries) — κ structure could depend on intent distribution", False),
        ("Stratified sample mismatch with Thakur 2025 — ours balanced, theirs natural-proportional", False),
        ("Gemini coverage failures (17-24% on TREC RAG 2024) — selection bias on the small valid subset", True),
        ("DSV4 Flash dropped — OpenRouter free-tier 429 throttling; we have 1 DeepSeek judge, no within-DS pair", False),
        ("Missing-data accounting — None=0 in aggregates; valid-only means re-cluster Gemini judges as strict-mid", False),
        ("Mediation analysis assumes linear additivity — non-linear interactions (provider × reasoning) untested", False),
        ("No within-corpus human-rater study — §7 NIST validation substitutes; reviewers preferring institution-specific human anchor will not find one", True),
    ], size=16, line_spacing=1.25)
    add_footer(s, f"P4 — {version_label}")
    add_slide_number(s, n, total)


def slide_reproducibility(p, n, total):
    s = standard_content_slide(p, "Reproducibility — single-command replication")
    add_textbox(s, 0.7, 1.4, 12.5, 0.5, "Within-corpus replication (~$18, ~5.5 h):",
                size=18, bold=True, color=COLOR_TITLE)
    add_textbox(s, 0.7, 1.95, 12.5, 0.5,
                "py eval_llm_judge.py --collection <yours> --judge-preset p4-frontier",
                size=15, font="Consolas", color=COLOR_TEXT)
    add_textbox(s, 0.7, 2.4, 12.5, 0.5,
                "py eval_llm_judge.py --collection <yours> --judge-preset p4-supplement-openweight",
                size=15, font="Consolas", color=COLOR_TEXT)
    add_textbox(s, 0.7, 3.1, 12.5, 0.5, "External-validation replication (~$35, ~7.5 h):",
                size=18, bold=True, color=COLOR_TITLE)
    add_textbox(s, 0.7, 3.65, 12.5, 0.5,
                "py validate_against_trec.py --corpus trec-rag-2024 --judge-preset p4-frontier --max-pairs 537",
                size=15, font="Consolas", color=COLOR_TEXT)
    add_textbox(s, 0.7, 4.10, 12.5, 0.5,
                "py validate_against_trec.py --corpus trec-rag-2024 --judge-preset p4-supplement-openweight --max-pairs 537",
                size=15, font="Consolas", color=COLOR_TEXT)
    add_textbox(s, 0.7, 4.55, 12.5, 0.5,
                "py validate_against_trec.py --analyze trec-rag-2024",
                size=15, font="Consolas", color=COLOR_TEXT)
    add_textbox(s, 0.7, 5.4, 12.5, 1.5,
                "Model versions pinned 2026-04-25: claude-opus-4-7, claude-sonnet-4-6, gpt-5.5, gpt-4o-2024-08-06, "
                "gemini-3.1-pro-preview, gemini-2.5-pro-2025-06-17, deepseek-v4-pro, qwen3.6-plus, gemma-4-26b-a4b-it. "
                "Target SIGIR Artifact Badging (Functional / Reusable).",
                size=13, color=COLOR_GREY)
    add_footer(s, "P4 — LONG")
    add_slide_number(s, n, total)


def slide_future_work(p, n, total):
    s = standard_content_slide(p, "Open questions for follow-up")
    add_bullets(s, 0.7, 1.4, 12.5, 5.0, [
        ("(1) Cross-organization Qwen ↔ Gemma 4 convergence (κ=0.80) lacks complete mechanistic explanation — "
         "shared OSS research community? distillation lineage? compute regime? bounded-output prior?", True),
        ("(2) κ structure under non-bounded output rubrics (graded relevance with longer scales, free-text rationales) "
         "may invert the within-family bounded-ordinal finding (C2)", False),
        ("(3) Non-linear interactions in §6 mediation analysis (provider × reasoning) untested", False),
        ("(4) Apples-to-apples replication of Thakur 2025's 537-pair subset pending receipt of their curated sample", False),
        ("(5) TREC-COVID and LLMJudge benchmark validations deferred — TOIS journal-version extension", False),
        ("(6) Cross-corpus mean κ on 4+ public benchmarks (vs 1 in current draft) for stronger generalization claim", True),
    ], size=17, line_spacing=1.3)
    add_footer(s, "P4 — LONG")
    add_slide_number(s, n, total)


# ----------------------- Build SHORT deck -----------------------
def build_short() -> Presentation:
    p = new_pres()
    slides = [
        ("title", slide_title),
        ("motivation", slide_motivation),
        ("contributions", slide_contribution_overview),
        ("judge slate", slide_judge_slate),
        ("pipeline", slide_pipeline_diagram),
        ("metrics", slide_per_judge_metrics),
        ("kappa matrix", slide_kappa_matrix_9j),
        ("calibration topology", slide_calibration_topology),
        ("C1+C2", slide_C1),
        ("C3+C4", slide_C3),
        ("mechanism", slide_mechanism_summary),
        ("ext val", slide_external_validation_results),
        ("coverage", slide_coverage_finding),
        ("takeaway", slide_takeaway),
        ("Q&A", slide_qa),
    ]
    total = len(slides)
    for n, (label, fn) in enumerate(slides, start=1):
        fn(p, n, total, version_label="SHORT")
    return p


# ----------------------- Build LONG deck -----------------------
def build_long() -> Presentation:
    p = new_pres()
    slides = [
        # Front matter
        ("title", lambda p, n, t: slide_title(p, n, t, version_label="LONG")),
        ("outline", slide_outline_long),
        ("motivation", lambda p, n, t: slide_motivation(p, n, t, version_label="LONG")),
        ("contributions overview", lambda p, n, t: slide_contribution_overview(p, n, t, version_label="LONG")),
        # Related work
        ("related work", slide_related_work),
        # Methodology
        ("corpus", slide_corpus_setup),
        ("judge slate", lambda p, n, t: slide_judge_slate(p, n, t, version_label="LONG")),
        ("rubric", slide_rubric),
        ("pipeline", slide_pipeline_diagram_long),
        # Within-corpus results
        ("metrics", lambda p, n, t: slide_per_judge_metrics(p, n, t, version_label="LONG")),
        ("calibration tiers", slide_calibration_tiers),
        ("kappa matrix 9j", lambda p, n, t: slide_kappa_matrix_9j(p, n, t, version_label="LONG")),
        ("calibration topology", lambda p, n, t: slide_calibration_topology(p, n, t, version_label="LONG")),
        # Findings
        ("C1", lambda p, n, t: slide_C1(p, n, t, version_label="LONG")),
        ("C2", lambda p, n, t: slide_C2(p, n, t, version_label="LONG")),
        ("C3", lambda p, n, t: slide_C3(p, n, t, version_label="LONG")),
        ("C4", lambda p, n, t: slide_C4(p, n, t, version_label="LONG")),
        # Mechanism (6 slides)
        ("mechanism summary", lambda p, n, t: slide_mechanism_summary(p, n, t, version_label="LONG")),
        ("pair confusion top-4 highest", slide_pair_confusion_top),
        ("pair confusion bottom-4 lowest", slide_pair_confusion_bottom),
        ("regression panels A+B", slide_regression_AB),
        ("regression panels C+D", slide_regression_CD),
        ("tokenizer refutation", slide_tokenizer_refutation),
        # External validation
        ("ext val setup", lambda p, n, t: slide_external_validation_setup(p, n, t, version_label="LONG")),
        ("ext val results", lambda p, n, t: slide_external_validation_results(p, n, t, version_label="LONG")),
        ("coverage", lambda p, n, t: slide_coverage_finding(p, n, t, version_label="LONG")),
        ("Thakur comparison", slide_thakur_comparison),
        # Closing
        ("limitations", lambda p, n, t: slide_limitations(p, n, t, version_label="LONG")),
        ("reproducibility", slide_reproducibility),
        ("future work", slide_future_work),
        ("takeaway", lambda p, n, t: slide_takeaway(p, n, t, version_label="LONG")),
        ("Q&A", lambda p, n, t: slide_qa(p, n, t, version_label="LONG")),
    ]
    total = len(slides)
    for n, (label, fn) in enumerate(slides, start=1):
        fn(p, n, total)
    return p


def main() -> int:
    print(f"Figures dir: {FIGS}")
    print(f"Output dir: {ROOT}")
    print()

    print("Building P4-SHORT.pptx...")
    short = build_short()
    short.save(str(OUT_SHORT))
    print(f"  ✓ saved {OUT_SHORT.name} ({OUT_SHORT.stat().st_size//1024} KB, {len(short.slides)} slides)")
    print()

    print("Building P4-LONG.pptx...")
    long = build_long()
    long.save(str(OUT_LONG))
    print(f"  ✓ saved {OUT_LONG.name} ({OUT_LONG.stat().st_size//1024} KB, {len(long.slides)} slides)")
    print()

    print("Done. Re-run anytime to regenerate from the latest figures + numbers.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
